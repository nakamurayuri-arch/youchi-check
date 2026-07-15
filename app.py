# -*- coding: utf-8 -*-
"""
用地チェック（緯度経度だけ）— 1ファイル完結・全項目
必要ライブラリ: streamlit, requests, Pillow
"""
import re, math, time, os
from io import BytesIO
import requests
import json
import streamlit.components.v1 as components
import streamlit as st

UA = {"User-Agent": "youchi-check/1.0"}
T = 25

st.set_page_config(page_title="用地チェック（緯度経度）", layout="centered")
st.title("用地チェック（緯度経度だけ）")
st.caption("緯度経度を入れるだけで用地の各項目を自動判定します。会社情報は使いません。")

def get_key():
    try:
        if "REINFOLIB_KEY" in st.secrets:
            return st.secrets["REINFOLIB_KEY"]
    except Exception:
        pass
    return os.environ.get("REINFOLIB_KEY")
REINFOLIB_KEY = get_key()

def parse_coord(s):
    s = (s or "").strip()
    m = re.search(r"(-?\d+\.\d+)\s*[, ]\s*(-?\d+\.\d+)", s)
    if m:
        return float(m.group(1)), float(m.group(2))
    m = re.search(r'(\d+)[°](\d+)[\'′](\d+\.?\d*)"?\s*N[, ]*\s*(\d+)[°](\d+)[\'′](\d+\.?\d*)"?\s*E', s, re.I)
    if m:
        lat = int(m.group(1)) + int(m.group(2))/60 + float(m.group(3))/3600
        lon = int(m.group(4)) + int(m.group(5))/60 + float(m.group(6))/3600
        return round(lat, 7), round(lon, 7)
    return None

def parse_one(s):
    """1つの座標値を10進 or 度分秒から解釈して float(度) を返す。N/E=+, S/W=-。"""
    if s is None:
        return None
    s = str(s).strip().replace("　", " ")
    if s == "":
        return None
    # 10進（例 33.944370 / -33.94 / 33.94N）
    m = re.fullmatch(r"\s*([+-]?\d+(?:\.\d+)?)\s*([NSEWnsew北南東西]?)\s*", s)
    if m:
        v = float(m.group(1))
        if m.group(2) in ("S", "s", "W", "w", "南", "西"):
            v = -abs(v)
        return v
    # 度分秒（例 33°56'39.7"N / N33 56 39.7 / 33 56 39.7）
    m = re.search(r"([NSEWnsew北南東西])?\s*(\d+)[°\s]+(\d+)['′\s]+(\d+(?:\.\d+)?)[\"″]?\s*([NSEWnsew北南東西])?", s)
    if m:
        d, mnt, sec = int(m.group(2)), int(m.group(3)), float(m.group(4))
        v = d + mnt/60 + sec/3600
        hemi = (m.group(1) or m.group(5) or "").upper()
        if hemi in ("S", "W", "南", "西"):
            v = -v
        return round(v, 7)
    # 度分（例 33 56.66）
    m = re.fullmatch(r"\s*([NSEWnsew北南東西])?\s*(\d+)[°\s]+(\d+(?:\.\d+)?)['′]?\s*([NSEWnsew北南東西])?\s*", s)
    if m:
        v = int(m.group(2)) + float(m.group(3))/60
        hemi = (m.group(1) or m.group(4) or "").upper()
        if hemi in ("S", "W", "南", "西"):
            v = -v
        return round(v, 7)
    return None

import xml.etree.ElementTree as _ET

def parse_kml(text):
    """KML文字列 → [{'name','type':'polygon'|'point','coords':[(lon,lat),...]}]。名前空間非依存。"""
    def local(tag): return tag.split('}')[-1]
    def coords(t):
        pts=[]
        for tok in (t or '').replace('\n',' ').replace('\t',' ').split():
            a=tok.split(',')
            if len(a)>=2:
                try: pts.append((float(a[0]), float(a[1])))
                except Exception: pass
        return pts
    root=_ET.fromstring(text)
    out=[]
    for pm in root.iter():
        if local(pm.tag)!='Placemark': continue
        name=None
        for ch in pm:
            if local(ch.tag)=='name': name=(ch.text or '').strip()
        poly_done=False
        for poly in pm.iter():
            if local(poly.tag)=='Polygon':
                for c in poly.iter():
                    if local(c.tag)=='coordinates' and c.text:
                        p=coords(c.text)
                        if p: out.append({'name':name,'type':'polygon','coords':p}); poly_done=True
                        break
                break
        if not poly_done:
            for pt in pm.iter():
                if local(pt.tag)=='Point':
                    for c in pt.iter():
                        if local(c.tag)=='coordinates' and c.text:
                            p=coords(c.text)
                            if p: out.append({'name':name,'type':'point','coords':p})
                            break
    return out

def kml_centroid(coords):
    if not coords: return None
    xs=[c[0] for c in coords]; ys=[c[1] for c in coords]
    return (sum(ys)/len(ys), sum(xs)/len(xs))  # (lat, lon)

LEAFLET_TEMPLATE = """
<!DOCTYPE html><html><head>
<link rel="stylesheet" href="https://unpkg.com/leaflet@1.9.4/dist/leaflet.css"/>
<script src="https://unpkg.com/leaflet@1.9.4/dist/leaflet.js"></script>
<style>#map{height:460px;border-radius:8px;}</style></head>
<body><div id="map"></div><script>
var lat=__LAT__, lon=__LON__, zoom=__ZOOM__, poly=__POLY__, zoning=__ZONING__, farm=__FARM__;
var map=L.map('map').setView([lat,lon],zoom);
L.tileLayer('https://cyberjapandata.gsi.go.jp/xyz/pale/{z}/{x}/{y}.png',{attribution:'国土地理院',maxZoom:18}).addTo(map);
function hz(u){return L.tileLayer('https://disaportaldata.gsi.go.jp/raster/'+u+'/{z}/{x}/{y}.png',{opacity:0.6,maxZoom:17});}
var ov={
 '洪水(想定最大規模)':hz('01_flood_l2_shinsuishin'),
 '津波':hz('04_tsunami_newlegend_data'),
 '高潮':hz('03_hightide_l2_shinsuishin_data'),
 '土砂(土石流)':hz('05_dosekiryukeikaikuiki_data'),
 '土砂(急傾斜)':hz('05_kyukeishakeikaikuiki_data'),
 '土砂(地すべり)':hz('05_jisuberikeikaikuiki_data')
};
ov['洪水(想定最大規模)'].addTo(map);
if(zoning){
  var zl=L.geoJSON(zoning,{style:function(f){
    var c=(f.properties&&f.properties.area_classification_ja)||'';
    var col=(c.indexOf('調整')>=0)?'#CB0F4B':((c.indexOf('市街化区域')>=0)?'#2E6D8A':'#999');
    return {color:col,weight:1,fillOpacity:0.25};
  },onEachFeature:function(f,l){var c=(f.properties&&f.properties.area_classification_ja)||'';l.bindPopup('区域区分: '+c);}});
  ov['市街化区域区分(赤=調整/青=市街化)']=zl;
}
if(farm){
  var fl=L.geoJSON(farm,{style:function(f){
    var ln=(f.properties&&f.properties.LAYER_NO);
    var col=(ln==6)?'#1F8A4C':'#C8A02E';
    return {color:col,weight:1,fillOpacity:0.25};
  },onEachFeature:function(f,l){var ln=(f.properties&&f.properties.LAYER_NO);l.bindPopup(ln==6?'農用地区域(青地)':'農業地域(白地含む)');}});
  ov['農地 青地(緑)/白地(黄)']=fl;
}
L.control.layers(null,ov,{collapsed:false}).addTo(map);
L.marker([lat,lon]).addTo(map);
if(poly){L.polygon(poly,{color:'#20376C',weight:2,fillOpacity:0.05}).addTo(map);}
</script></body></html>
"""

def leaflet_hazard_html(lat, lon, poly_lonlat=None, zoom=16, zoning=None, farm=None):
    poly_js = json.dumps([[y, x] for (x, y) in poly_lonlat]) if poly_lonlat else "null"
    return (LEAFLET_TEMPLATE
            .replace("__LAT__", repr(float(lat)))
            .replace("__LON__", repr(float(lon)))
            .replace("__ZOOM__", str(int(zoom)))
            .replace("__POLY__", poly_js)
            .replace("__ZONING__", json.dumps(zoning) if zoning else "null")
            .replace("__FARM__", json.dumps(farm) if farm else "null"))

def reinfolib_geojson_around(endpoint, lat, lon, z=15, span=1):
    """地点を含むタイル＋周辺(span)タイルのフィーチャを集めてFeatureCollection化。地図描画用。"""
    if not REINFOLIB_KEY:
        return {"type": "FeatureCollection", "features": []}
    xf, yf = deg2tile(lat, lon, z); xc, yc = int(xf), int(yf)
    feats = []
    for dx in range(-span, span+1):
        for dy in range(-span, span+1):
            try:
                for f in _reinfolib_tile(endpoint, xc+dx, yc+dy, z):
                    g = f.get("geometry", {})
                    if g.get("type") in ("Polygon", "MultiPolygon"):
                        feats.append({"type": "Feature",
                                      "properties": f.get("properties", {}),
                                      "geometry": g})
            except Exception:
                pass
    return {"type": "FeatureCollection", "features": feats}

def a12_geojson_around(lat, lon, addr, dd=0.012):
    """A12(青地/白地)を地点周辺(±dd度)だけGeoJSON化。LAYER_NO(6=青地,5=白地)をpropsに。"""
    feats = []
    for code, pname in _pref_candidates(lat, lon, addr):
        try:
            shps = download_a12(code)
        except Exception:
            continue
        import shapefile
        got = False
        for shp in shps:
            try:
                r = shapefile.Reader(shp, encoding="cp932")
            except Exception:
                r = shapefile.Reader(shp)
            fields = [f[0] for f in r.fields[1:]]
            li = fields.index("LAYER_NO") if "LAYER_NO" in fields else None
            for sr in r.iterShapeRecords():
                bb = sr.shape.bbox
                if bb[2] < lon-dd or bb[0] > lon+dd or bb[3] < lat-dd or bb[1] > lat+dd:
                    continue
                try:
                    geom = sr.shape.__geo_interface__
                except Exception:
                    continue
                ln = None
                if li is not None:
                    try: ln = int(sr.record[li])
                    except Exception: ln = sr.record[li]
                feats.append({"type": "Feature", "geometry": geom, "properties": {"LAYER_NO": ln}})
                got = True
        if got:
            break
    return {"type": "FeatureCollection", "features": feats}

def sample_points(coords, grid=3):
    """区画ポリゴン(coords=[(lon,lat)]) → 頂点+辺中点+重心+内部グリッド点。"""
    n = len(coords)
    if n == 0:
        return []
    pts = list(coords)
    for i in range(n):
        x1, y1 = coords[i]; x2, y2 = coords[(i+1) % n]
        pts.append(((x1+x2)/2, (y1+y2)/2))
    cx = sum(p[0] for p in coords)/n; cy = sum(p[1] for p in coords)/n
    pts.append((cx, cy))
    xs = [p[0] for p in coords]; ys = [p[1] for p in coords]
    xmin, xmax, ymin, ymax = min(xs), max(xs), min(ys), max(ys)
    def _pip(x, y, ring):
        inside = False; m = len(ring); j = m-1
        for i in range(m):
            xi, yi = ring[i]; xj, yj = ring[j]
            if ((yi > y) != (yj > y)) and (x < (xj-xi)*(y-yi)/(yj-yi)+xi):
                inside = not inside
            j = i
        return inside
    for gx in range(1, grid):
        for gy in range(1, grid):
            x = xmin+(xmax-xmin)*gx/grid; y = ymin+(ymax-ymin)*gy/grid
            if _pip(x, y, coords):
                pts.append((x, y))
    return list({(round(x, 7), round(y, 7)) for (x, y) in pts})

def _reinfolib_tile(endpoint, xt, yt, z):
    url = f"https://www.reinfolib.mlit.go.jp/ex-api/external/{endpoint}?response_format=geojson&z={z}&x={xt}&y={yt}"
    r = requests.get(url, headers={**UA, "Ocp-Apim-Subscription-Key": REINFOLIB_KEY}, timeout=T)
    if r.status_code != 200:
        return []
    gj = r.json()
    return gj.get("features", []) if isinstance(gj, dict) else []

def face_reinfolib(endpoint, pts, z):
    """各サンプル点で該当した全ポリゴンのpropsリストを返す（タイルキャッシュ）。"""
    cache = {}; out = []
    for (lon, lat) in pts:
        xf, yf = deg2tile(lat, lon, z); key = (int(xf), int(yf))
        if key not in cache:
            try:
                cache[key] = _reinfolib_tile(endpoint, key[0], key[1], z)
            except Exception:
                cache[key] = []
        matched = [f.get("properties", {}) for f in cache[key] if point_in_geom(lon, lat, f.get("geometry", {}))]
        out.append(matched)
    return out

def _coverage(n_hit, n):
    if n == 0: return "不明"
    if n_hit == 0: return "非該当"
    if n_hit == n: return "区画ほぼ全域が該当"
    return "区画の一部が該当"

def face_flood(pts):
    lists = face_reinfolib("XKT026", pts, 15)
    ranks = []
    for ml in lists:
        rs = [int(p.get("A31a_205") or 0) for p in ml if p.get("A31a_205")]
        if rs: ranks.append(max(rs))
    n_hit = len(ranks); n = len(pts)
    mx = max(ranks) if ranks else None
    return {"cover": _coverage(n_hit, n), "n_hit": n_hit, "n": n,
            "max_depth": FLOOD_RANK.get(mx) if mx else None, "max_rank": mx}

def face_kubun(pts):
    lists = face_reinfolib("XKT001", pts, 15)
    labels = [_kubun_from_feats(ml) for ml in lists]
    n_chosei = sum(1 for l in labels if l == "市街化調整区域")
    n_shigai = sum(1 for l in labels if l == "市街化区域")
    return {"n_chosei": n_chosei, "n_shigai": n_shigai, "n": len(pts),
            "cover_chosei": _coverage(n_chosei, len(pts))}

def face_a12(pts, addr):
    st_counts = {"青地": 0, "白地": 0, "非農地": 0}
    for (lon, lat) in pts:
        try:
            s = a12_status(lat, lon, addr).get("status")
        except Exception:
            s = None
        if s in st_counts: st_counts[s] += 1
    n = len(pts)
    return {"counts": st_counts, "n": n,
            "cover_aochi": _coverage(st_counts["青地"], n)}


def haversine(a, b, c, d):
    R = 6371000.0; r = math.radians
    h = math.sin(r(c-a)/2)**2 + math.cos(r(a))*math.cos(r(c))*math.sin(r(d-b)/2)**2
    return 2*R*math.asin(math.sqrt(h))

def deg2tile(lat, lon, z):
    n = 2**z
    return (lon+180)/360*n, (1-math.asinh(math.tan(math.radians(lat)))/math.pi)/2*n

def _pip_ring(x, y, ring):
    inside = False; n = len(ring); j = n-1
    for i in range(n):
        xi, yi = ring[i][0], ring[i][1]; xj, yj = ring[j][0], ring[j][1]
        if ((yi > y) != (yj > y)) and (x < (xj-xi)*(y-yi)/(yj-yi) + xi):
            inside = not inside
        j = i
    return inside

def point_in_geom(x, y, geom):
    t = geom.get("type"); c = geom.get("coordinates")
    if t == "Polygon" and c:
        if not _pip_ring(x, y, c[0]):
            return False
        return not any(_pip_ring(x, y, h) for h in c[1:])
    if t == "MultiPolygon" and c:
        for poly in c:
            if poly and _pip_ring(x, y, poly[0]) and not any(_pip_ring(x, y, h) for h in poly[1:]):
                return True
    return False

def elevation(lat, lon):
    r = requests.get("https://cyberjapandata2.gsi.go.jp/general/dem/scripts/getelevation.php",
                     params={"lon": lon, "lat": lat, "outtype": "JSON"}, headers=UA, timeout=T)
    return float(r.json().get("elevation"))

def slope(lat, lon, step=20.0):
    dlat = step/111000; dlon = step/(111000*math.cos(math.radians(lat)))
    z0 = elevation(lat, lon); mx = 0.0
    for la, lo in [(lat+dlat, lon), (lat-dlat, lon), (lat, lon+dlon), (lat, lon-dlon)]:
        try:
            mx = max(mx, math.degrees(math.atan2(abs(elevation(la, lo)-z0), step)))
        except Exception:
            pass
    return round(z0, 1), round(mx, 1)

HAZ = {
    "洪水浸水(想定最大規模)": "https://disaportaldata.gsi.go.jp/raster/01_flood_l2_shinsuishin_data/{z}/{x}/{y}.png",
    "津波浸水想定": "https://disaportaldata.gsi.go.jp/raster/04_tsunami_newlegend_data/{z}/{x}/{y}.png",
    "高潮浸水想定": "https://disaportaldata.gsi.go.jp/raster/03_hightide_l2_shinsuishin_data/{z}/{x}/{y}.png",
    "土砂(土石流)": "https://disaportaldata.gsi.go.jp/raster/05_dosekiryukeikaikuiki_data/{z}/{x}/{y}.png",
    "土砂(急傾斜)": "https://disaportaldata.gsi.go.jp/raster/05_kyukeishakeikaikuiki_data/{z}/{x}/{y}.png",
    "土砂(地すべり)": "https://disaportaldata.gsi.go.jp/raster/05_jisuberikeikaikuiki_data/{z}/{x}/{y}.png",
}
DEPTH_LEGEND = [
    ((247, 245, 169), "0〜0.5m未満"),
    ((255, 216, 192), "0.5〜3m未満"),
    ((255, 183, 183), "3〜5m未満"),
    ((255, 145, 145), "5〜10m未満"),
    ((242, 133, 201), "10〜20m未満"),
    ((220, 122, 220), "20m以上"),
]
DOSHA_LEGEND = [
    ((255, 237, 74), "警戒区域（イエロー）"),
    ((203, 15, 75), "特別警戒区域（レッド）"),
    ((250, 122, 122), "特別警戒区域（レッド）"),
]
DEPTH_LAYERS = {"洪水浸水(想定最大規模)", "津波浸水想定", "高潮浸水想定"}

def classify_color(rgb, palette):
    best, bd = None, 1e9
    for col, label in palette:
        d = sum((a-b)**2 for a, b in zip(rgb, col))
        if d < bd:
            bd, best = d, label
    return best

def hazard(lat, lon, z=16):
    """タイルで該当/非該当のみ判定（色→深さ推定は廃止。深さは洪水のみXKT026で正確取得）。"""
    from PIL import Image
    xf, yf = deg2tile(lat, lon, z); xt, yt = int(xf), int(yf)
    out = {}
    for name, tmpl in HAZ.items():
        try:
            r = requests.get(tmpl.format(z=z, x=xt, y=yt), headers=UA, timeout=T)
            if r.status_code != 200:
                out[name] = False; continue
            img = Image.open(BytesIO(r.content)).convert("RGBA")
            px = min(255, int((xf-xt)*256)); py = min(255, int((yf-yt)*256))
            d = img.getpixel((px, py))
            out[name] = bool(d[3] > 0 and (d[0]+d[1]+d[2] > 0))
        except Exception:
            out[name] = None
        time.sleep(0.1)
    return out

# 洪水浸水深ランク（国土数値情報A31 / reinfolib XKT026 の A31a_205）
FLOOD_RANK = {1:"0〜0.5m未満",2:"0.5〜1m未満",3:"1〜2m未満",4:"2〜3m未満",
              5:"3〜4m未満",6:"4〜5m未満",7:"5〜10m未満",8:"10〜20m未満",9:"20m以上"}

def reinfolib_flood(lat, lon, z=15):
    """XKT026(洪水/想定最大規模)。該当ポリゴン全ての(河川,ランク)を集め最大を採用。診断付き。"""
    if not REINFOLIB_KEY:
        return None
    xf, yf = deg2tile(lat, lon, z); xt, yt = int(xf), int(yf)
    url = f"https://www.reinfolib.mlit.go.jp/ex-api/external/XKT026?response_format=geojson&z={z}&x={xt}&y={yt}"
    try:
        r = requests.get(url, headers={**UA, "Ocp-Apim-Subscription-Key": REINFOLIB_KEY}, timeout=T)
        if r.status_code in (204, 404):
            return {"hit": False, "matches": []}
        if r.status_code != 200:
            return {"err": f"HTTP {r.status_code}"}
        gj = r.json()
    except Exception as e:
        return {"err": str(e)}
    matches = []
    total = len(gj.get("features", []) if isinstance(gj, dict) else [])
    for f in (gj.get("features", []) if isinstance(gj, dict) else []):
        if point_in_geom(lon, lat, f.get("geometry", {})):
            p = f.get("properties", {})
            try:
                rk = int(p.get("A31a_205") or 0)
            except Exception:
                rk = 0
            matches.append({"river": p.get("A31a_202"), "rank": rk,
                            "depth": FLOOD_RANK.get(rk, "?"), "props_keys": list(p.keys())})
    if not matches:
        return {"hit": False, "matches": [], "tile_features": total}
    mr = max(m["rank"] for m in matches)
    top = max(matches, key=lambda m: m["rank"])
    return {"hit": True, "rank": mr, "depth": FLOOD_RANK.get(mr, "不明"),
            "river": top.get("river"), "matches": matches, "tile_features": total}

def reinfolib_inundation(endpoint, lat, lon, z=14):
    """XKT027(高潮)/XKT028(津波)：浸水深は文字列区分。属性から深さ文字列を拾い最深を返す。"""
    if not REINFOLIB_KEY:
        return None
    xf, yf = deg2tile(lat, lon, z); xt, yt = int(xf), int(yf)
    url = f"https://www.reinfolib.mlit.go.jp/ex-api/external/{endpoint}?response_format=geojson&z={z}&x={xt}&y={yt}"
    try:
        r = requests.get(url, headers={**UA, "Ocp-Apim-Subscription-Key": REINFOLIB_KEY}, timeout=T)
        if r.status_code in (204, 404):
            return {"hit": False}
        if r.status_code != 200:
            return {"err": f"HTTP {r.status_code}"}
        gj = r.json()
    except Exception as e:
        return {"err": str(e)}
    hit_any = False; depths = []
    for f in (gj.get("features", []) if isinstance(gj, dict) else []):
        if point_in_geom(lon, lat, f.get("geometry", {})):
            hit_any = True
            for v in (f.get("properties", {}) or {}).values():
                if isinstance(v, str) and re.search(r"\d", v) and "m" in v and ("以上" in v or "未満" in v or "以下" in v):
                    depths.append(v)
    if not hit_any:
        return {"hit": False}
    if depths:
        def _dk(s):
            n = re.findall(r"\d+(?:\.\d+)?", s); return float(n[0]) if n else 0.0
        return {"hit": True, "depth": max(depths, key=_dk)}
    return {"hit": True, "depth": None}

REINFOLIB = [
    ("市街化区域/調整区域", "XKT001", ["kubun_id", "kubun", "区域区分"]),
    ("用途地域", "XKT002", ["use_area_ja"]),
    ("自然公園地域", "XKT019", ["ptext", "park_name", "公園名", "thema"]),
    ("地すべり防止区域", "XKT021", ["name", "区域名"]),
    ("急傾斜地崩壊危険区域", "XKT022", ["name", "区域名"]),
]

def reinfolib_hit(lat, lon, endpoint, z=15):
    xf, yf = deg2tile(lat, lon, z); xt, yt = int(xf), int(yf)
    url = f"https://www.reinfolib.mlit.go.jp/ex-api/external/{endpoint}?response_format=geojson&z={z}&x={xt}&y={yt}"
    r = requests.get(url, headers={**UA, "Ocp-Apim-Subscription-Key": REINFOLIB_KEY}, timeout=T)
    if r.status_code != 200:
        return {"err": f"HTTP {r.status_code}"}
    gj = r.json()
    feats = gj.get("features", []) if isinstance(gj, dict) else []
    matched = [f.get("properties", {}) for f in feats if point_in_geom(lon, lat, f.get("geometry", {}))]
    if matched:
        return {"hit": True, "props": matched[0], "all_props": matched}
    return {"hit": False}

def _kubun_from_feats(all_props):
    """XKT001の該当フィーチャ群から 市街化区域/市街化調整区域/非線引き を判定。"""
    vals = [str(p.get("area_classification_ja") or "") for p in (all_props or [])]
    if any("市街化調整区域" in v for v in vals):
        return "市街化調整区域"
    if any("市街化区域" in v for v in vals):
        return "市街化区域"
    if any("都市計画区域" in v for v in vals):
        return "非線引き都市計画区域（区域区分なし）"
    return None

def label_from(props, keys):
    for k in keys:
        if k in props and props[k] not in (None, ""):
            return str(props[k])
    for v in props.values():
        if isinstance(v, str) and v.strip():
            return v
    return "該当"


# ---------- 農地：青地/白地（国土数値情報 A12 を都道府県別に取得→点内包） ----------
PREF_CODE = {
 "北海道":"01","青森県":"02","岩手県":"03","宮城県":"04","秋田県":"05","山形県":"06","福島県":"07",
 "茨城県":"08","栃木県":"09","群馬県":"10","埼玉県":"11","千葉県":"12","東京都":"13","神奈川県":"14",
 "新潟県":"15","富山県":"16","石川県":"17","福井県":"18","山梨県":"19","長野県":"20","岐阜県":"21",
 "静岡県":"22","愛知県":"23","三重県":"24","滋賀県":"25","京都府":"26","大阪府":"27","兵庫県":"28",
 "奈良県":"29","和歌山県":"30","鳥取県":"31","島根県":"32","岡山県":"33","広島県":"34","山口県":"35",
 "徳島県":"36","香川県":"37","愛媛県":"38","高知県":"39","福岡県":"40","佐賀県":"41","長崎県":"42",
 "熊本県":"43","大分県":"44","宮崎県":"45","鹿児島県":"46","沖縄県":"47",
}

def pref_code_from_addr(addr):
    if not addr:
        return None
    for name, code in PREF_CODE.items():
        if name in addr:
            return code, name
    return None

@st.cache_data(show_spinner=False)
def download_a12(code):
    import io, zipfile, tempfile, glob
    url = f"https://nlftp.mlit.go.jp/ksj/gml/data/A12/A12-15/A12-15_{code}_GML.zip"
    r = requests.get(url, headers=UA, timeout=120)
    r.raise_for_status()
    d = tempfile.mkdtemp()
    zipfile.ZipFile(io.BytesIO(r.content)).extractall(d)
    shps = glob.glob(d + "/**/*.shp", recursive=True)
    if not shps:
        raise RuntimeError("A12のSHPが見つかりません")
    return shps

# 近隣県（県境の取り違え対策）: 主要な隣接関係のみ簡易に
NEIGHBORS = {
 "25":["26","21","24","18","23"], "26":["25","27","28","29","24","18"],
 "27":["26","28","29","30"], "28":["27","26","31","33","36","24"],
 "13":["11","12","14","19"], "14":["13","11","19","22"],
}

def _pref_from_latlon(lat, lon):
    """緯度経度→都道府県名（逆ジオコード）。住所が無くても県を特定するため。"""
    try:
        js = requests.get("https://msearch.gsi.go.jp/address-search/AddressSearch",
                          params={"q": f"{lat},{lon}"}, headers=UA, timeout=T).json()
    except Exception:
        js = None
    # GSIの緯度経度→住所は別APIなのでreverseを使う
    try:
        j = requests.get("https://nominatim.openstreetmap.org/reverse",
                         params={"format": "json", "accept-language": "ja", "lat": lat, "lon": lon},
                         headers=UA, timeout=T).json()
        for name in PREF_CODE:
            if name in (j.get("display_name") or "") or name in str(j.get("address", {})):
                return name
    except Exception:
        pass
    return None

def _pref_candidates(lat, lon, addr):
    cands = []
    seen = set()
    def add(code):
        if code and code not in seen:
            nm = [k for k, v in PREF_CODE.items() if v == code]
            if nm:
                seen.add(code); cands.append((code, nm[0]))
    pc = pref_code_from_addr(addr)
    if pc:
        add(pc[0])
        for nb in NEIGHBORS.get(pc[0], []):
            add(nb)
    # 住所で取れない/不足なら、緯度経度の逆ジオコードで県を特定
    if not cands:
        nm = _pref_from_latlon(lat, lon)
        if nm:
            code = PREF_CODE[nm]
            add(code)
            for nb in NEIGHBORS.get(code, []):
                add(nb)
    return cands

def _rings_from_shape(shape):
    pts = shape.points
    parts = list(shape.parts) + [len(pts)]
    return [pts[parts[i]:parts[i+1]] for i in range(len(parts)-1)]

def _point_in_rings(lon, lat, rings):
    # even-odd rule：内包するリング数が奇数なら内部（外周/穴を自動処理）
    cnt = 0
    for ring in rings:
        if len(ring) >= 3 and _pip_ring(lon, lat, ring):
            cnt += 1
    return (cnt % 2) == 1

def _hit_in_shps(shps, lat, lon):
    import shapefile
    for shp in shps:
        try:
            r = shapefile.Reader(shp, encoding="cp932")
        except Exception:
            r = shapefile.Reader(shp)
        fields = [f[0] for f in r.fields[1:]]
        for sr in r.iterShapeRecords():
            bb = sr.shape.bbox
            if not (bb[0] <= lon <= bb[2] and bb[1] <= lat <= bb[3]):
                continue
            if _point_in_rings(lon, lat, _rings_from_shape(sr.shape)):
                return {"attrs": dict(zip(fields, list(sr.record))), "file": shp.split("/")[-1]}
    return None

def _diag_shps(shps, lat, lon):
    """SHPを読み、件数・全体bbox・点がbbox内か等の診断を返す。"""
    import shapefile
    total = 0; inbbox = 0
    gxmin = gymin = 1e18; gxmax = gymax = -1e18
    sample_fields = []
    for shp in shps:
        try:
            r = shapefile.Reader(shp, encoding="cp932")
        except Exception:
            r = shapefile.Reader(shp)
        sample_fields = [f[0] for f in r.fields[1:]]
        b = r.bbox  # [xmin,ymin,xmax,ymax]
        gxmin = min(gxmin, b[0]); gymin = min(gymin, b[1])
        gxmax = max(gxmax, b[2]); gymax = max(gymax, b[3])
        for sh in r.iterShapes():
            total += 1
            bb = sh.bbox
            if bb[0] <= lon <= bb[2] and bb[1] <= lat <= bb[3]:
                inbbox += 1
    return {"files": [s.split("/")[-1] for s in shps], "polygons": total,
            "data_bbox": [round(gxmin,4), round(gymin,4), round(gxmax,4), round(gymax,4)],
            "point": [round(lon,6), round(lat,6)],
            "point_in_data_bbox": (gxmin <= lon <= gxmax and gymin <= lat <= gymax),
            "bbox_hit_polygons": inbbox, "fields": sample_fields}

def _a12_layers_hit(shps, lat, lon):
    """点を含むポリゴンの LAYER_NO 集合（6=農用地区域/5=農業地域）と代表属性を返す。"""
    import shapefile
    layers = set(); sample = None; sfile = None
    for shp in shps:
        try:
            r = shapefile.Reader(shp, encoding="cp932")
        except Exception:
            r = shapefile.Reader(shp)
        fields = [f[0] for f in r.fields[1:]]
        li = fields.index("LAYER_NO") if "LAYER_NO" in fields else None
        for sr in r.iterShapeRecords():
            bb = sr.shape.bbox
            if not (bb[0] <= lon <= bb[2] and bb[1] <= lat <= bb[3]):
                continue
            if _point_in_rings(lon, lat, _rings_from_shape(sr.shape)):
                ln = None
                if li is not None:
                    try: ln = int(sr.record[li])
                    except Exception: ln = sr.record[li]
                layers.add(ln)
                if sample is None:
                    sample = dict(zip(fields, list(sr.record))); sfile = shp.split("/")[-1]
    return layers, sample, sfile

def a12_status(lat, lon, addr):
    cands = _pref_candidates(lat, lon, addr)
    if not cands:
        return {"err": "都道府県を特定できず"}
    tried = []; diag = None
    for code, pname in cands:
        try:
            shps = download_a12(code)
        except Exception as e:
            tried.append(f"{pname}(コード{code}):取得失敗 {e}"); continue
        if diag is None:
            try:
                diag = {"pref": pname, "code": code, **_diag_shps(shps, lat, lon)}
            except Exception:
                pass
        layers, sample, sfile = _a12_layers_hit(shps, lat, lon)
        if 6 in layers:
            return {"status": "青地", "pref": pname, "attrs": sample, "file": sfile, "diag": diag}
        if layers:  # 5（農業地域）にのみ該当
            return {"status": "白地", "pref": pname, "attrs": sample, "file": sfile, "diag": diag}
    return {"status": "非農地", "pref": cands[0][1], "note": "／".join(tried), "diag": diag}

OVP = "https://overpass-api.de/api/interpreter"
ROADCLS = {"trunk": "国道相当", "primary": "国道/主要地方道", "secondary": "県道相当",
           "tertiary": "主要市町村道", "unclassified": "一般道", "residential": "生活道路",
           "service": "私道/構内", "track": "農道/林道"}

def ovp(q):
    return requests.get(OVP, params={"data": q}, headers=UA, timeout=50).json()

def nearest_road(lat, lon):
    js = ovp(f"[out:json][timeout:40];way(around:250,{lat},{lon})[highway];out tags geom;")
    best = None
    for e in js.get("elements", []):
        dmin = min((haversine(lat, lon, p["lat"], p["lon"]) for p in e.get("geometry", [])), default=1e12)
        c = {"d": round(dmin), "cls": ROADCLS.get(e["tags"].get("highway"), e["tags"].get("highway")),
             "name": e["tags"].get("name") or e["tags"].get("ref") or ""}
        if best is None or c["d"] < best["d"]:
            best = c
    return best

def _osm_building(lat, lon):
    js = ovp(f"[out:json][timeout:40];way(around:400,{lat},{lon})[building];out geom;")
    dmin = 1e12
    for e in js.get("elements", []):
        for p in e.get("geometry", []):
            dmin = min(dmin, haversine(lat, lon, p["lat"], p["lon"]))
    return round(dmin) if dmin < 1e12 else None

def _extent_to_latlon(px, py, X, Y, Z, extent=4096, y_down=False):
    fx = px/extent; fy = py/extent
    gx = X + fx
    gy = (Y + (1-fy)) if not y_down else (Y + fy)
    n = 2**Z
    lon = gx/n*360 - 180
    lat = math.degrees(math.atan(math.sinh(math.pi*(1 - 2*gy/n))))
    return lat, lon

def _iter_rings(geom):
    t = geom.get("type"); c = geom.get("coordinates")
    if t == "Polygon":
        for ring in c: yield ring
    elif t == "MultiPolygon":
        for poly in c:
            for ring in poly: yield ring
    elif t == "LineString":
        yield c
    elif t == "Point":
        yield [c]

def _gsi_building(lat, lon, z=16):
    """国土地理院ベクトルタイル(建物)から最寄り建物距離。OSMの穴（工業地/埋立地/農村）を補う。"""
    try:
        import mapbox_vector_tile as mvt
    except Exception:
        return None
    xf, yf = deg2tile(lat, lon, z); X, Y = int(xf), int(yf)
    best = 1e12
    for dx in (0, -1, 1):
        for dy in (0, -1, 1):
            tx, ty = X+dx, Y+dy
            try:
                r = requests.get(f"https://cyberjapandata.gsi.go.jp/xyz/experimental_bvmap/{z}/{tx}/{ty}.pbf",
                                 headers=UA, timeout=T)
                if r.status_code != 200 or not r.content:
                    continue
                data = mvt.decode(r.content)
            except Exception:
                continue
            layer = data.get("building") or data.get("建物")
            if not layer:
                continue
            extent = layer.get("extent", 4096)
            for f in layer.get("features", []):
                for ring in _iter_rings(f.get("geometry", {})):
                    for pt in ring:
                        blat, blon = _extent_to_latlon(pt[0], pt[1], tx, ty, z, extent)
                        d = haversine(lat, lon, blat, blon)
                        if d < best:
                            best = d
    return round(best) if best < 1e12 else None

def nearest_building(lat, lon):
    vals = []
    try:
        v = _osm_building(lat, lon)
        if v is not None: vals.append(v)
    except Exception:
        pass
    try:
        v = _gsi_building(lat, lon)
        if v is not None: vals.append(v)
    except Exception:
        pass
    return min(vals) if vals else None

def _ovp_geom(q, timeout=60):
    r = requests.get(OVP, params={"data": q}, headers=UA, timeout=timeout)
    return r.json()

def coast_dist(lat, lon):
    """海（海岸線・海面）までの最短距離＋診断。海岸線と海面を別クエリにして切り分ける。"""
    diag = {"n_coast": 0, "n_water": 0, "err": None, "dist": None}
    dmin = 1e12
    try:
        js = _ovp_geom(f"[out:json][timeout:55];way(around:15000,{lat},{lon})[natural=coastline];out geom;")
        for el in js.get("elements", []):
            diag["n_coast"] += 1
            for p in el.get("geometry", []) or []:
                if p and "lat" in p:
                    dmin = min(dmin, haversine(lat, lon, p["lat"], p["lon"]))
    except Exception as ex:
        diag["err"] = f"coastline:{ex}"
    try:
        q = (f"[out:json][timeout:55];("
             f"way(around:8000,{lat},{lon})[natural=water][water~\"sea|bay|strait|lagoon|harbour\"];"
             f"relation(around:8000,{lat},{lon})[natural=water][water~\"sea|bay|strait|lagoon|harbour\"];"
             f"way(around:8000,{lat},{lon})[natural=bay];"
             f"way(around:6000,{lat},{lon})[landuse=harbour];"
             f"way(around:6000,{lat},{lon})[harbour=yes];"
             f");out geom;")
        js = _ovp_geom(q)
        for el in js.get("elements", []):
            diag["n_water"] += 1
            for p in el.get("geometry", []) or []:
                if p and "lat" in p:
                    dmin = min(dmin, haversine(lat, lon, p["lat"], p["lon"]))
            for mem in el.get("members", []) or []:
                for p in mem.get("geometry", []) or []:
                    if p and "lat" in p:
                        dmin = min(dmin, haversine(lat, lon, p["lat"], p["lon"]))
    except Exception as ex:
        diag["err"] = (diag["err"] or "") + f" water:{ex}"
    diag["dist"] = round(dmin) if dmin < 1e12 else None
    return diag

def revgeo(lat, lon):
    try:
        return requests.get("https://nominatim.openstreetmap.org/reverse",
                            params={"format": "json", "accept-language": "ja", "lat": lat, "lon": lon},
                            headers=UA, timeout=T).json().get("display_name")
    except Exception:
        return None

if not REINFOLIB_KEY:
    st.warning("市街化調整・用途・自然公園・地すべり・急傾斜を自動判定するには、無料のreinfolib APIキーが必要です（申請 → Streamlitの Secrets に REINFOLIB_KEY を設定）。未設定の間はリンク表示になります。", icon="🔑")

# ========== PowerPoint出力（テンプレ＋合成画像） ==========
import os as _os
GSI_PHOTO = "https://cyberjapandata.gsi.go.jp/xyz/seamlessphoto/{z}/{x}/{y}.jpg"
GSI_PALE = "https://cyberjapandata.gsi.go.jp/xyz/pale/{z}/{x}/{y}.png"
TEMPLATE_PATH = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "dd_template.pptx")

def _fetch_img(url, rgba=False):
    try:
        from PIL import Image
        r = requests.get(url, headers=UA, timeout=T)
        if r.status_code != 200:
            return None
        im = Image.open(BytesIO(r.content))
        return im.convert("RGBA") if rgba else im.convert("RGB")
    except Exception:
        return None

def _draw_geom_on(draw, geom, topx, fill, outline=None):
    t = geom.get("type"); c = geom.get("coordinates")
    polys = [c] if t == "Polygon" else (c if t == "MultiPolygon" else [])
    for poly in polys:
        if poly and len(poly[0]) >= 3:
            pts = [topx(x, y) for (x, y) in poly[0]]
            draw.polygon(pts, fill=fill)
            if outline:
                draw.line(pts + [pts[0]], fill=outline, width=1)

def compose_map(lat, lon, z=16, ntiles=3, overlay_url=None, zoning=None, farm=None, poly=None, label=None, base="photo"):
    from PIL import Image, ImageDraw
    xf, yf = deg2tile(lat, lon, z); xc, yc = int(xf), int(yf); half = ntiles // 2
    W = H = ntiles * 256
    canvas = Image.new("RGB", (W, H), (210, 210, 210))
    for dx in range(-half, half + 1):
        for dy in range(-half, half + 1):
            _burl = GSI_PHOTO if base == "photo" else GSI_PALE
            t = _fetch_img(_burl.format(z=z, x=xc + dx, y=yc + dy))
            if t:
                canvas.paste(t.resize((256, 256)), ((dx + half) * 256, (dy + half) * 256))
    if overlay_url:
        layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
        for dx in range(-half, half + 1):
            for dy in range(-half, half + 1):
                t = _fetch_img(overlay_url.format(z=z, x=xc + dx, y=yc + dy), rgba=True)
                if t:
                    t = t.resize((256, 256))
                    layer.paste(t, ((dx + half) * 256, (dy + half) * 256), t)
        _a = layer.getchannel("A").point(lambda a: int(a * 0.6))
        layer.putalpha(_a)
        canvas = Image.alpha_composite(canvas.convert("RGBA"), layer).convert("RGB")
    draw = ImageDraw.Draw(canvas, "RGBA")
    def topx(lo, la):
        x, y = deg2tile(la, lo, z); return ((x - (xc - half)) * 256, (y - (yc - half)) * 256)
    if zoning:
        for f in zoning.get("features", []):
            cc = (f.get("properties", {}).get("area_classification_ja") or "")
            if "調整" in cc:
                fill, line = (203, 15, 75, 64), (203, 15, 75, 255)
            elif "市街化区域" in cc:
                fill, line = (46, 109, 138, 64), (46, 109, 138, 255)
            else:
                fill, line = (150, 150, 150, 48), (150, 150, 150, 200)
            _draw_geom_on(draw, f.get("geometry", {}), topx, fill, outline=line)
    if farm:
        for f in farm.get("features", []):
            ln = f.get("properties", {}).get("LAYER_NO")
            if ln == 6:
                fill, line = (31, 138, 76, 64), (31, 138, 76, 255)
            else:
                fill, line = (200, 160, 46, 64), (200, 160, 46, 255)
            _draw_geom_on(draw, f.get("geometry", {}), topx, fill, outline=line)
    if poly:
        pts = [topx(lo, la) for (lo, la) in poly]
        draw.line(pts + [pts[0]], fill=(203, 15, 75, 255), width=4)
    cx, cy = topx(lon, lat)
    draw.ellipse([cx - 7, cy - 7, cx + 7, cy + 7], fill=(32, 55, 108, 255), outline=(255, 255, 255, 255), width=2)
    if label:
        draw.rectangle([0, 0, 8 + 7 * len(label), 22], fill=(255, 255, 255, 210))
        draw.text((5, 5), label, fill=(20, 20, 20))
    out = BytesIO(); canvas.save(out, "PNG"); return out.getvalue()

def build_report_pptx(rows, subtitle, images):
    import copy as _copy
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.text.text import _Paragraph
    RED = RGBColor(0xC0, 0x00, 0x00)
    def setc(cell, text, bold=False, red=False):
        tf = cell.text_frame; lines = str(text).split("\n")
        for _pp in tf.paragraphs[1:]:
            _pp._p.getparent().remove(_pp._p)
        p0 = tf.paragraphs[0]; r0 = p0.runs[0] if p0.runs else p0.add_run()
        r0.text = lines[0]; r0.font.bold = bold
        if red: r0.font.color.rgb = RED
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
        for extra in lines[1:]:
            newp = _copy.deepcopy(p0._p); p0._p.getparent().append(newp)
            pp = _Paragraph(newp, p0._parent)
            for rr in pp.runs[1:]:
                rr._r.getparent().remove(rr._r)
            pp.runs[0].text = extra; pp.runs[0].font.bold = bold
            if red: pp.runs[0].font.color.rgb = RED
    prs = Presentation(TEMPLATE_PATH); s = prs.slides[0]
    for sh in s.shapes:
        if sh.name == "タイトル 1" and sh.has_text_frame:
            sh.text_frame.paragraphs[0].runs[0].text = "用地調査結果"
        if sh.name == "コンテンツ プレースホルダー 4" and sh.has_text_frame:
            tf = sh.text_frame
            (tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()).text = subtitle
    tbl = [sh for sh in s.shapes if sh.has_table][0].table
    tbl.columns[0].width = Inches(2.5); tbl.columns[1].width = Inches(3.0)
    nd = len(tbl.rows) - 1
    for i, (k, v, hit) in enumerate(rows):
        if i + 1 <= nd:
            setc(tbl.cell(i + 1, 0), k, bold=True)
            setc(tbl.cell(i + 1, 1), v, bold=bool(hit), red=bool(hit))
    for ri in range(nd, len(rows), -1):
        tr = tbl.rows[ri]._tr; tr.getparent().remove(tr)
    remove = {"図 31", "図 9", "図 13", "Rectangle 15", "Straight Connector 17",
              "Straight Connector 18", "5-Point Star 23", "矢印: 下 21", "正方形/長方形 22", "矢印: 下 24"}
    for sh in list(s.shapes):
        if sh.name in remove:
            sh._element.getparent().remove(sh._element)
    if images:
        s.shapes.add_picture(BytesIO(images[0]), Inches(6.75), Inches(1.45), Inches(6.15), Inches(3.3))
        rest = images[1:5]; n = len(rest)
        if n:
            gap = 0.15; total = 6.15
            w = min(3.0, (total - gap * (n - 1)) / n); h = min(1.9, w * 0.66); T = 4.95
            for i, img in enumerate(rest):
                s.shapes.add_picture(BytesIO(img), Inches(6.75 + i * (w + gap)), Inches(T), Inches(w), Inches(h))
    out = BytesIO(); prs.save(out); return out.getvalue()

addr_in = st.text_input("住所（記録用。緯度経度が空のときはここから判定）",
                        placeholder="例: 滋賀県野洲市堤字ノ爪740-1")
_c1, _c2 = st.columns(2)
with _c1:
    lat_in = st.text_input("緯度", placeholder="例: 33.944370 / 33°56'39.7\"N")
with _c2:
    lon_in = st.text_input("経度", placeholder="例: 130.807654 / 130°48'27.6\"E")
kml_file = st.file_uploader("区画KML（任意・Googleマイマップ等から書き出し。区画があればその重心で判定＋地図に表示）",
                            type=["kml"])
if st.button("▶ チェックする", type="primary"):
    kml_feats = []
    if kml_file is not None:
        try:
            kml_feats = parse_kml(kml_file.getvalue().decode("utf-8", "ignore"))
        except Exception as _ke:
            st.warning(f"KMLの読取に失敗しました: {_ke}")
    lat = lon = None
    if kml_feats:
        _tgt = next((f for f in kml_feats if f["type"] == "polygon"), kml_feats[0])
        _cen = kml_centroid(_tgt["coords"])
        if _cen:
            lat, lon = _cen
            st.caption(f"KMLの区画『{_tgt.get('name') or '無題'}』の重心で判定： {lat:.6f}, {lon:.6f}")
    if lat is None or lon is None:
        lat = parse_one(lat_in)
        lon = parse_one(lon_in)
    if lat is not None and lon is not None:
        if addr_in.strip() and not kml_feats:
            st.caption(f"判定は緯度経度を使用（住所『{addr_in.strip()}』は記録用）")
    elif addr_in.strip():
        g = geocode(addr_in.strip())
        if not g:
            st.error("住所から場所を特定できませんでした。番地まで入れるか、緯度経度を入力してください。")
            st.stop()
        lat, lon, _ = g
        st.caption(f"住所を緯度経度に変換： {lat:.6f}, {lon:.6f}")
    else:
        st.error("住所か緯度経度のどちらかを入力してください。")
        st.stop()
    if kml_feats:
        try:
            import pydeck as pdk
            polys = [{"polygon": [[x, y] for (x, y) in f["coords"]], "name": f.get("name") or ""}
                     for f in kml_feats if f["type"] == "polygon"]
            layers = []
            if polys:
                layers.append(pdk.Layer("PolygonLayer", polys, get_polygon="polygon",
                                        get_fill_color=[203, 15, 75, 50], get_line_color=[203, 15, 75],
                                        line_width_min_pixels=2, pickable=True))
            layers.append(pdk.Layer("ScatterplotLayer", [{"position": [lon, lat]}], get_position="position",
                                    get_fill_color=[32, 55, 108], get_radius=15, radius_min_pixels=5))
            st.subheader("区画マップ")
            st.pydeck_chart(pdk.Deck(layers=layers, map_style=None,
                            initial_view_state=pdk.ViewState(latitude=lat, longitude=lon, zoom=15),
                            tooltip={"text": "{name}"}))
        except Exception as _me:
            st.caption(f"（区画マップの描画をスキップ: {_me}）")
        _poly = next((f for f in kml_feats if f["type"] == "polygon"), None)
        if _poly:
            with st.spinner("区画の面判定中…"):
                _pts = sample_points(_poly["coords"], grid=3)
                st.subheader(f"区画（面）判定 ｜ サンプル{len(_pts)}点")
                if REINFOLIB_KEY:
                    ff = face_flood(_pts)
                    if ff.get("max_rank"):
                        st.write(f"- 洪水（面）： **{ff['cover']}**（区画内 最大 {ff['max_depth']}／{ff['n_hit']}/{ff['n']}点該当）")
                    else:
                        st.write(f"- 洪水（面）： **{ff['cover']}**")
                    fk = face_kubun(_pts)
                    if fk["n_chosei"] > 0:
                        st.write(f"- 市街化調整（面）： **{fk['cover_chosei']}**（調整区域 {fk['n_chosei']}/{fk['n']}点）")
                    else:
                        st.write(f"- 市街化調整（面）： **区画内に市街化調整区域なし**（市街化区域 {fk['n_shigai']}/{fk['n']}点）")
                else:
                    st.write("- 洪水・市街化調整（面）： reinfolib APIキー設定後に面判定")
                fa = face_a12(_pts, addr_in)
                _c = fa["counts"]
                st.write(f"- 青地/白地（面）： 青地 {_c['青地']}／白地 {_c['白地']}／非農地 {_c['非農地']}（全{fa['n']}点）→ 青地は{fa['cover_aochi']}")
                st.caption("面判定は区画の頂点・辺・内部のサンプル点による近似です。境界付近は公式マップで確認してください。")
    with st.spinner("判定中…（10〜40秒）"):
        addr = revgeo(lat, lon)
        try: elev, slp = slope(lat, lon)
        except Exception: elev = slp = None
        haz = hazard(lat, lon)
        rein = {}
        if REINFOLIB_KEY:
            for name, ep, keys in REINFOLIB:
                try:
                    rein[name] = reinfolib_hit(lat, lon, ep)
                except Exception as e:
                    rein[name] = {"err": str(e)}
                time.sleep(0.3)
        try: road = nearest_road(lat, lon)
        except Exception: road = None
        try: bldg = nearest_building(lat, lon)
        except Exception: bldg = None
        try: cd = coast_dist(lat, lon)
        except Exception as _e: cd = {"dist": None, "err": str(_e), "n_coast": 0, "n_water": 0}
        coast = cd.get("dist")

    st.subheader("基本")
    st.write({"緯度経度": f"{lat:.6f}, {lon:.6f}", "住所(推定)": addr or "取得できず",
              "標高": f"{elev} m" if elev is not None else "取得できず",
              "傾斜(推定)": f"{slp}°" if slp is not None else "取得できず"})

    st.subheader("ハザード")
    _poly_ll = None
    if kml_feats:
        _pf = next((f for f in kml_feats if f["type"] == "polygon"), None)
        if _pf:
            _poly_ll = _pf["coords"]
    with st.spinner("地図レイヤ（区域区分・農地）を準備中…"):
        _zoning = reinfolib_geojson_around("XKT001", lat, lon, z=15, span=1) if REINFOLIB_KEY else None
        try:
            _farm = a12_geojson_around(lat, lon, addr_in)
        except Exception:
            _farm = None
    try:
        components.html(leaflet_hazard_html(lat, lon, _poly_ll, zoom=16, zoning=_zoning, farm=_farm), height=470)
        st.caption("アプリ内地図（右上で 洪水/津波/高潮/土砂、市街化区域区分、農地 を切替）。外部サイトを開かないため位置ずれは起きません。出典：ハザードマップポータルサイト／国土地理院、不動産情報ライブラリ、国土数値情報A12。")
    except Exception as _le:
        st.caption(f"（地図の埋め込みをスキップ: {_le}）")
    fl  = reinfolib_flood(lat, lon) if REINFOLIB_KEY else None
    tsu = reinfolib_inundation("XKT028", lat, lon, z=14) if REINFOLIB_KEY else None
    hig = reinfolib_inundation("XKT027", lat, lon, z=14) if REINFOLIB_KEY else None

    def _line_flood(k):
        if isinstance(fl, dict) and not fl.get("err"):
            if fl.get("hit"):
                rv = f"（{fl.get('river','')}）" if fl.get("river") else ""
                return f"- {k}： **⚠ 該当（{fl.get('depth','')}／最大ランク{fl.get('rank')}）**{rv} 〈XKT026〉"
            return f"- {k}： **○ 非該当** 〈XKT026〉"
        return None

    def _line_str(k, res, ep):
        if isinstance(res, dict) and not res.get("err"):
            if res.get("hit"):
                d = f"（{res['depth']}）" if res.get("depth") else "（深さ区分の記載なし）"
                return f"- {k}： **⚠ 該当{d}** 〈{ep}〉"
            return f"- {k}： **○ 非該当** 〈{ep}〉"
        return None

    for k, v in haz.items():
        line = None
        if k.startswith("洪水"):
            line = _line_flood(k)
        elif k.startswith("津波"):
            line = _line_str(k, tsu, "XKT028")
        elif k.startswith("高潮"):
            line = _line_str(k, hig, "XKT027")
        if line is None:
            if v is True:
                extra = "（深さは公式マップで確認）" if k.startswith(("洪水", "津波", "高潮")) else ""
                line = f"- {k}： **⚠ 該当**{extra}"
            elif v is False:
                line = f"- {k}： **○ 非該当**"
            else:
                line = f"- {k}： 要確認"
        st.write(line)
    if isinstance(fl, dict) and fl.get("hit"):
        with st.expander("洪水の内訳（XKT026が返した該当ポリゴン）"):
            st.json({"tile_features": fl.get("tile_features"), "matches": fl.get("matches")})
    for nm, res in [("洪水", fl), ("津波", tsu), ("高潮", hig)]:
        if isinstance(res, dict) and res.get("err"):
            st.caption(f"（{nm}の属性取得でエラー: {res['err']}。タイル判定にフォールバック）")

    st.subheader("許認可・区域")
    if REINFOLIB_KEY:
        for name, ep, keys in REINFOLIB:
            r = rein.get(name, {})
            if name == "市街化区域/調整区域":
                if r.get("err"):
                    st.write(f"- 市街化調整区域： 取得エラー（{r['err']}）")
                elif r.get("hit"):
                    lab = _kubun_from_feats(r.get("all_props", []))
                    if lab == "市街化調整区域":
                        st.write("- 市街化調整区域： **⚠ 該当（市街化調整区域）**")
                    elif lab == "市街化区域":
                        st.write("- 市街化調整区域： **○ 非該当（市街化区域）**")
                    elif lab:
                        st.write(f"- 市街化調整区域： **○ 非該当**（{lab}）")
                    else:
                        st.write("- 市街化調整区域： 都市計画区域内だが区分不明 → 下の属性を確認")
                    with st.expander("XKT001の該当属性（全ポリゴン）"):
                        st.json(r.get("all_props", []))
                else:
                    st.write("- 市街化調整区域： **○ 非該当**（都市計画区域外の可能性）")
                continue
            if r.get("err"):
                st.write(f"- {name}： 取得エラー（{r['err']}）")
            elif r.get("hit"):
                st.write(f"- {name}： **⚠ 該当**（{label_from(r.get('props', {}), keys)}）")
            else:
                st.write(f"- {name}： **○ 非該当**")
    else:
        st.markdown("- 市街化調整・用途・自然公園・地すべり・急傾斜： reinfolib APIキー設定後に自動判定（上のアプリ内地図でも確認できます）")

    st.subheader("周辺")
    bt = "周辺に建物なし" if bldg is None else (f"{bldg} m ⚠100m未満" if bldg < 100 else f"{bldg} m")
    ct = "海岸を検出できず(内陸の可能性)" if coast is None else (f"{coast} m ⚠500m未満(重塩害注意)" if coast < 500 else f"{coast} m")
    st.write(f"- 最寄り道路： **{(road['cls']+' '+road['name']+' '+str(road['d'])+'m') if road else '取得できず'}**")
    st.write(f"- 最寄り建物(住宅目安)： **{bt}**")
    st.write(f"- 海岸まで(重塩害)： **{ct}**")
    with st.expander("海岸判定の診断"):
        st.json(cd)

    st.subheader("農地（青地/白地：国土数値情報A12・2015年度）")
    try:
        a12 = a12_status(lat, lon, ((addr_in or "") + " " + (addr or "")).strip())
    except Exception as e:
        a12 = {"err": str(e)}
    if a12.get("err"):
        st.write(f"- 農地区分： 取得エラー（{a12['err']}）")
    else:
        stt = a12.get("status")
        if stt == "青地":
            st.write(f"- 農地区分： **⚠ 青地（農用地区域内）**（判定県：{a12.get('pref','')}）")
            st.caption("→ 農地転用には原則、農振除外（農振法13条2）が先行して必要。")
        elif stt == "白地":
            st.write(f"- 農地区分： **△ 白地（農業地域内・農用地区域外）**（判定県：{a12.get('pref','')}）")
            st.caption("→ 農地転用の対象（第何種は農業委員会確認）。")
        else:
            st.write(f"- 農地区分： **○ 非農地／農業地域外**（判定県：{a12.get('pref','')}）")
        if a12.get("attrs"):
            with st.expander("A12の属性（生データ）"):
                st.json({"file": a12.get("file"), "attrs": a12.get("attrs")})
        if a12.get("note"):
            st.warning(f"データ取得の注意: {a12['note']}")
    with st.expander("A12 判定の診断"):
        st.json(a12.get("diag") or {"注意": "診断情報なし"})
    st.caption("※A12は2015年度データ。農用地区域（青地）は参考表示で精度保証なし。第何種は農業委員会、地目は登記簿で確認。")

    st.subheader("その他（データ無し→リンク確認）")
    st.markdown(
        f"- 農地の地目・地番・第何種： [農地ナビ](https://map.maff.go.jp/)\n"
        f"- 砂防指定地： [国土数値情報](https://nlftp.mlit.go.jp/ksj/)／都道府県の砂防GIS\n"
        f"- 埋蔵文化財包蔵地： [文化財総覧WebGIS](https://heritagemap.nabunken.go.jp/)（全国統一の可否データ無し・自治体教委）\n"
        f"- 位置： [Googleマップ](https://www.google.com/maps?q={lat},{lon})"
    )
    st.info("公開データからの一次確認です。最終確定は各自治体・現地・公式マップで。")
    if REINFOLIB_KEY:
        st.caption("このサービスは、国土交通省不動産情報ライブラリのAPI機能を使用していますが、提供情報の最新性、正確性、完全性等が保証されたものではありません。")

    # ===== レポート用データ集約（PowerPoint出力用） =====
    def _r_kubun():
        if not REINFOLIB_KEY: return ("reinfolibキー未設定", False)
        r = rein.get("市街化区域/調整区域", {})
        if r.get("hit"):
            lab = _kubun_from_feats(r.get("all_props", []))
            if lab == "市街化調整区域": return ("該当（市街化調整区域）", True)
            if lab == "市街化区域": return ("非該当（市街化区域）", False)
            if lab: return ("非該当（" + lab + "）", False)
            return ("都市計画区域内（区分不明）", False)
        return ("非該当（都市計画区域外の可能性）", False)
    def _r_youto():
        if not REINFOLIB_KEY: return ("reinfolibキー未設定", False)
        r = rein.get("用途地域", {})
        return (label_from(r.get("props", {}), ["use_area_ja"]), True) if r.get("hit") else ("非該当", False)
    def _r_nochi():
        stt = a12.get("status")
        if stt == "青地": return ("青地（農用地区域内）※農振除外が必要", True)
        if stt == "白地": return ("白地（農業地域内・農用地区域外）", False)
        if a12.get("err"): return ("取得エラー", False)
        return ("非農地／農業地域外", False)
    def _r_flood():
        if isinstance(fl, dict) and not fl.get("err"):
            if fl.get("hit"):
                rv = "／" + str(fl.get("river")) if fl.get("river") else ""
                return ("該当（" + str(fl.get("depth", "")) + rv + "）", True)
            return ("非該当", False)
        v = haz.get("洪水浸水(想定最大規模)")
        return ("該当" if v else "非該当", bool(v))
    def _r_tsuhig():
        parts = []; hit = False
        for nm, res, tl in [("津波", tsu, "津波浸水想定"), ("高潮", hig, "高潮浸水想定")]:
            if isinstance(res, dict) and not res.get("err"):
                if res.get("hit"):
                    parts.append(nm + "該当（" + (res.get("depth") or "深さ区分なし") + "）"); hit = True
                else:
                    parts.append(nm + "非該当")
            else:
                v = haz.get(tl); parts.append(nm + ("該当" if v else "非該当")); hit = hit or bool(v)
        return ("／".join(parts), hit)
    def _r_dosha():
        hit = any(haz.get(k) for k in ["土砂(土石流)", "土砂(急傾斜)", "土砂(地すべり)"])
        return ("該当（警戒/特別警戒）" if hit else "非該当", hit)
    def _r_kuiki():
        if not REINFOLIB_KEY: return ("reinfolibキー未設定", False)
        hits = [nm for nm, r in [("自然公園", rein.get("自然公園地域", {})),
                                 ("地すべり", rein.get("地すべり防止区域", {})),
                                 ("急傾斜", rein.get("急傾斜地崩壊危険区域", {}))] if r.get("hit")]
        return ("該当：" + "・".join(hits) if hits else "非該当", bool(hits))
    _coast = cd.get("dist") if isinstance(cd, dict) else None
    _rows = [
        ("所在地", (addr_in or addr or ""), False),
        ("緯度経度", f"{lat:.6f}, {lon:.6f}", False),
        ("市街化調整区域",) + _r_kubun(),
        ("用途地域",) + _r_youto(),
        ("農地区分（青地/白地）",) + _r_nochi(),
        ("ハザード：洪水",) + _r_flood(),
        ("ハザード：津波・高潮",) + _r_tsuhig(),
        ("ハザード：土砂",) + _r_dosha(),
        ("自然公園・地すべり・急傾斜",) + _r_kuiki(),
        ("標高・傾斜", (f"標高 {elev}m ／ 傾斜 {slp}°" if elev is not None else "取得できず"), False),
        ("最寄り道路／建物",
         "道路 " + ((road["cls"] + " " + str(road["d"]) + "m") if road else "—") + " ／ 建物 " + (str(bldg) + "m" if bldg is not None else "—"), False),
        ("海岸距離（重塩害）", (f"{_coast}m" if _coast is not None else "海岸なし（内陸）"), (_coast is not None and _coast < 500)),
    ]
    _ov = []
    if _r_flood()[1]:
        _ov.append(("洪水", HAZ["洪水浸水(想定最大規模)"], False, False))
    if _r_tsuhig()[1]:
        if (isinstance(tsu, dict) and tsu.get("hit")) or haz.get("津波浸水想定"):
            _ov.append(("津波", HAZ["津波浸水想定"], False, False))
        if (isinstance(hig, dict) and hig.get("hit")) or haz.get("高潮浸水想定"):
            _ov.append(("高潮", HAZ["高潮浸水想定"], False, False))
    if _r_dosha()[1]:
        for k in ["土砂(土石流)", "土砂(急傾斜)", "土砂(地すべり)"]:
            if haz.get(k):
                _ov.append(("土砂", HAZ[k], False, False)); break
    if _r_kubun()[1] and _zoning:
        _ov.append(("市街化調整", None, True, False))
    if a12.get("status") in ("青地", "白地") and _farm:
        _ov.append(("農地", None, False, True))
    st.session_state["report"] = {
        "rows": _rows, "subtitle": (addr_in or addr or f"{lat:.6f}, {lon:.6f}"),
        "lat": lat, "lon": lon, "poly": _poly_ll,
        "overlays": _ov[:4], "zoning": _zoning, "farm": _farm,
    }
    st.session_state.pop("pptx_bytes", None)


# ========== PowerPoint出力ボタン ==========
if st.session_state.get("report"):
    st.divider()
    st.subheader("PowerPoint出力（Sustechテンプレ）")
    if not _os.path.exists(TEMPLATE_PATH):
        st.info("テンプレート dd_template.pptx をGitHubリポジトリ（app.pyと同じ場所）に追加すると、この地点の調査結果を1枚もの資料として出力できます。")
    else:
        if st.button("📊 PowerPointを生成", type="primary"):
            rep = st.session_state["report"]
            with st.spinner("スライド生成中（航空写真・該当地図を取得）…"):
                try:
                    imgs = [compose_map(rep["lat"], rep["lon"], poly=rep["poly"], label="航空写真", base="photo")]
                    for (label, ov_url, use_z, use_f) in rep["overlays"]:
                        imgs.append(compose_map(rep["lat"], rep["lon"], overlay_url=ov_url,
                                    zoning=rep["zoning"] if use_z else None,
                                    farm=rep["farm"] if use_f else None,
                                    poly=rep["poly"], label=label, base="pale"))
                    st.session_state["pptx_bytes"] = build_report_pptx(rep["rows"], rep["subtitle"], imgs)
                except Exception as _ex:
                    st.error(f"生成に失敗しました: {_ex}")
        if st.session_state.get("pptx_bytes"):
            st.download_button("⬇ PPTXをダウンロード", st.session_state["pptx_bytes"],
                file_name="用地調査結果.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
            st.caption("該当項目は赤字太字、右に航空写真＋該当した地図（ハザード・市街化調整・農地）を掲載。ロゴ・タイトル書式はテンプレ準拠。")
